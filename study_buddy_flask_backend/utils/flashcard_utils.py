import csv
import uuid
import tempfile
from io import BytesIO, StringIO
from typing import List, Tuple, cast
from openpyxl import Workbook
from openpyxl.worksheet.worksheet import Worksheet
import genanki
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.shapes.placeholder import SlidePlaceholder
from pptx.shapes.autoshape import Shape
import re
import subprocess
import os
import sys
from PIL import Image
import cairosvg  # type: ignore
import math
from markdown import markdown
from weasyprint import HTML

# Add the parent directory to the Python path so we can import utils
current_dir = os.path.dirname(os.path.abspath(__file__))
parent_dir = os.path.dirname(current_dir)
if parent_dir not in sys.path:
    sys.path.insert(0, parent_dir)

from utils.gpt.shared import render_math_in_html, generate_ai_title


def export_flashcards_to_csv(flashcards: List[Tuple[str, str]]) -> BytesIO:
    """
    Accepts a list of (question, answer) tuples and returns a CSV buffer.
    """
    output = StringIO()
    writer = csv.writer(output)
    writer.writerow(["Question", "Answer"])
    for q, a in flashcards:
        writer.writerow([q, a])
    
    # Convert to BytesIO for Flask response
    csv_content = output.getvalue()
    output.close()
    
    bytes_output = BytesIO()
    bytes_output.write(csv_content.encode('utf-8'))
    bytes_output.seek(0)
    return bytes_output


def export_flashcards_to_xlsx(flashcards: List[Tuple[str, str]]) -> BytesIO:
    """
    Accepts a list of (question, answer) tuples and returns an Excel (.xlsx) buffer.
    """
    wb = Workbook()
    ws = cast(Worksheet, wb.active)
    ws.title = "Flashcards"
    ws.append(["Question", "Answer"])
    for q, a in flashcards:
        ws.append([q, a])
    output = BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def render_latex_to_png(latex_expr: str, display_mode: bool = False) -> str:
    """
    Renders LaTeX to a high-res PNG file using the existing KaTeX renderer (katex_renderer.js).
    Returns the path to the PNG file, or an empty string on error.
    Cleans up temp files after use.
    """
    import shutil
    current_dir = os.path.dirname(os.path.abspath(__file__))
    backend_dir = os.path.dirname(os.path.dirname(current_dir))
    katex_path = os.path.join(backend_dir, 'study_buddy_flask_backend', 'katex_renderer.js')
    svg_path = None
    png_path = None
    try:
        result = subprocess.run(
            ['node', katex_path],
            input=latex_expr,
            capture_output=True,
            text=True,
            check=True
        )
        svg = result.stdout.strip()
        if not svg or '<svg' not in svg:
            print(f"❌ KaTeX SVG output invalid or empty for: {latex_expr}")
            print(f"SVG output: {svg}")
            return ""
        # Save SVG and convert to high-res PNG
        with tempfile.NamedTemporaryFile(suffix=".svg", delete=False) as svg_file:
            svg_file.write(svg.encode('utf-8'))
            svg_path = svg_file.name
        png_path = svg_path.replace('.svg', '.png')
        try:
            # Use scale to increase resolution (default 1.0, try 2.0 for sharper images)
            cairosvg.svg2png(url=svg_path, write_to=png_path, scale=2.0)  # type: ignore
        except Exception as e:
            print(f"❌ SVG to PNG conversion failed for: {latex_expr}")
            print(f"SVG content: {svg}")
            print(f"Error: {e}")
            if svg_path and os.path.exists(svg_path):
                os.remove(svg_path)
            return ""
        if not os.path.exists(png_path) or os.path.getsize(png_path) == 0:
            print(f"❌ PNG file not created or empty for: {latex_expr}")
            if svg_path and os.path.exists(svg_path):
                os.remove(svg_path)
            return ""
        return png_path
    except subprocess.CalledProcessError as e:
        print(f"❌ KaTeX render error: {e}")
        print(f"stderr: {e.stderr}")
        return ""
    except Exception as e:
        print(f"❌ Unexpected error in render_latex_to_png: {e}")
        return ""
    finally:
        # Clean up SVG after PNG is created
        if svg_path and os.path.exists(svg_path):
            try:
                os.remove(svg_path)
            except Exception:
                pass


def add_text_and_math_to_shape(slide, text, left, top, width, height, line_spacing=0.4):
    """
    Adds text and math images to a slide, stacking them vertically and wrapping long text.
    """
    pattern = re.compile(r'(\$\$.*?\$\$|\$.*?\$|\\\[.*?\\\]|\\\(.*?\\\))', re.DOTALL)
    cursor = 0
    current_top = top
    line_height = height
    for match in pattern.finditer(text):
        start, end = match.span()
        # Add preceding text
        if start > cursor:
            plain = text[cursor:start].strip()
            if plain:
                textbox = slide.shapes.add_textbox(left, current_top, width, height)
                textbox.text = plain
                textbox.text_frame.word_wrap = True
                if textbox.text_frame.paragraphs:
                    textbox.text_frame.paragraphs[0].font.size = Pt(20)
                lines = max(1, math.ceil(len(plain) / 60))
                current_top += Inches(line_spacing) * lines
        # Add math image
        latex = match.group(0)
        # Remove delimiters
        if latex.startswith('$$') and latex.endswith('$$'):
            expr = latex[2:-2]
            display = True
        elif latex.startswith('$') and latex.endswith('$'):
            expr = latex[1:-1]
            display = False
        elif latex.startswith('\\[') and latex.endswith('\\]'):
            expr = latex[2:-2]
            display = True
        elif latex.startswith('\\(') and latex.endswith('\\)'):
            expr = latex[2:-2]
            display = False
        else:
            expr = latex
            display = False
        png_path = render_latex_to_png(expr, display)
        if png_path and os.path.exists(png_path) and os.path.getsize(png_path) > 0:
            try:
                # Center the image horizontally
                img_width = Inches(4) if display else Inches(2.5)
                img_left = left + (width - img_width) / 2
                print(f"Inserting PNG into slide: {png_path}")
                pic = slide.shapes.add_picture(png_path, img_left, current_top, img_width, Inches(0.7 if display else 0.5))
                current_top += pic.height + int(Inches(0.1))
                # Do NOT delete PNG after use (for debugging)
                # os.remove(png_path)
            except Exception as e:
                print(f"❌ Failed to add math image: {e}")
                # Fallback: insert LaTeX as text
                textbox = slide.shapes.add_textbox(left, current_top, width, height)
                textbox.text = latex
                textbox.text_frame.word_wrap = True
                if textbox.text_frame.paragraphs:
                    textbox.text_frame.paragraphs[0].font.size = Pt(20)
                lines = max(1, math.ceil(len(latex) / 60))
                current_top += Inches(line_spacing) * lines
        else:
            print(f"❌ Failed to render math: {expr}")
            # Fallback: insert LaTeX as text
            textbox = slide.shapes.add_textbox(left, current_top, width, height)
            textbox.text = latex
            textbox.text_frame.word_wrap = True
            if textbox.text_frame.paragraphs:
                textbox.text_frame.paragraphs[0].font.size = Pt(20)
            lines = max(1, math.ceil(len(latex) / 60))
            current_top += Inches(line_spacing) * lines
        cursor = end
    # Add any remaining text
    if cursor < len(text):
        plain = text[cursor:].strip()
        if plain:
            textbox = slide.shapes.add_textbox(left, current_top, width, height)
            textbox.text = plain
            textbox.text_frame.word_wrap = True
            if textbox.text_frame.paragraphs:
                textbox.text_frame.paragraphs[0].font.size = Pt(20)
            lines = max(1, math.ceil(len(plain) / 60))
            current_top += Inches(line_spacing) * lines


def export_flashcards_to_pptx(flashcards: List[Tuple[str, str]]) -> BytesIO:
    """
    Accepts a list of (question, answer) tuples and returns a PowerPoint (.pptx) buffer.
    Each question is a slide, followed by its answer as the next slide.
    Renders LaTeX math as images.
    """
    prs = Presentation()
    content_slide_layout = prs.slide_layouts[1]  # Title and Content
    for q, a in flashcards:
        # Question slide
        slide_q = prs.slides.add_slide(content_slide_layout)
        if slide_q.shapes.title is not None:
            title_q = slide_q.shapes.title
            title_q.text = "Question"
            if title_q.text_frame and title_q.text_frame.paragraphs:
                title_q.text_frame.paragraphs[0].font.size = Pt(28)
        # Add question text and math
        add_text_and_math_to_shape(slide_q, q, Inches(1), Inches(2), Inches(8), Inches(0.5), line_spacing=0.4)
        # Answer slide
        slide_a = prs.slides.add_slide(content_slide_layout)
        if slide_a.shapes.title is not None:
            title_a = slide_a.shapes.title
            title_a.text = "Answer"
            if title_a.text_frame and title_a.text_frame.paragraphs:
                title_a.text_frame.paragraphs[0].font.size = Pt(28)
        add_text_and_math_to_shape(slide_a, a, Inches(1), Inches(2), Inches(8), Inches(0.5), line_spacing=0.4)
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output


def format_latex(text: str) -> str:
    """
    Formats mathematical text to use display-style LaTeX if it contains common math symbols.
    Escapes backslashes to avoid formatting errors.
    """
    text = text.replace("\\", "\\\\")
    if any(sym in text for sym in ["\\int", "\\frac", "^", "_", "\\sum", "\\lim", "\\sqrt"]):
        return f"\\[{text}\\]"
    return text


def export_flashcards_to_apkg(flashcards: List[Tuple[str, str]], deck_name: str = "Study Buddy Deck") -> BytesIO:
    """
    Accepts a list of (question, answer) tuples and returns a .apkg buffer using genanki with LaTeX formatting.
    """
    deck_id = int(uuid.uuid4().int >> 64)  # Ensure it's within SQLite INTEGER range

    model = genanki.Model(
        1607392319,
        'Basic Model',
        fields=[
            {"name": "Question"},
            {"name": "Answer"}
        ],
        templates=[
            {
                "name": "Card 1",
                "qfmt": "{{Question}}",
                "afmt": "{{FrontSide}}<hr id='answer'>{{Answer}}"
            }
        ],
        latex_pre=r"""
\documentclass[12pt]{article}
\usepackage{amsmath}
\usepackage{amssymb}
\pagestyle{empty}
\begin{document}
""",
        latex_post=r"\end{document}"
    )

    deck = genanki.Deck(deck_id, deck_name)

    for question, answer in flashcards:
        try:
            note = genanki.Note(
                model=model,
                fields=[format_latex(question), format_latex(answer)]
            )
            deck.add_note(note)
        except Exception as e:
            print(f"⚠️ Skipping flashcard due to error: {e}")

    tmp_path = tempfile.NamedTemporaryFile(suffix=".apkg", delete=False).name
    genanki.Package(deck).write_to_file(tmp_path)

    with open(tmp_path, "rb") as f:
        output = BytesIO(f.read())
    output.seek(0)
    return output


def export_flashcards_to_pdf(flashcards: List[Tuple[str, str]]) -> BytesIO:
    """
    Accepts a list of (question, answer) tuples and returns a PDF buffer with multiple Q/A pairs per page, math rendered with KaTeX.
    """
    # Build HTML for all flashcards
    html_blocks = []
    for idx, (q, a) in enumerate(flashcards, 1):
        html_blocks.append(f"""
        <div class='flashcard'>
            <div class='question'><b>Q{idx}:</b> {render_math_in_html(q)}</div>
            <div class='answer'><b>A{idx}:</b> {render_math_in_html(a)}</div>
        </div>
        """)
    html_body = "\n".join(html_blocks)
    full_html = f"""
    <html>
    <head>
        <meta charset='utf-8'>
        <title>Flashcards</title>
        <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.css">
        <style>
            body {{ font-family: 'Georgia', serif; margin: 2em; color: #222; }}
            .flashcard {{ margin-bottom: 2em; padding: 1em; border-bottom: 1px solid #ccc; }}
            .question {{ font-size: 18px; margin-bottom: 0.5em; }}
            .answer {{ font-size: 17px; color: #2a6ebd; }}
        </style>
    </head>
    <body>
        <h1>Flashcards</h1>
        {html_body}
    </body>
    </html>
    """
    pdf_buffer = BytesIO()
    HTML(string=full_html).write_pdf(pdf_buffer)
    pdf_buffer.seek(0)
    return pdf_buffer


def export_flashcards_to_anki_tsv(flashcards):
    from io import BytesIO
    buffer = BytesIO()
    for q, a in flashcards:
        buffer.write(f"{q}\t{a}\n".encode("utf-8"))
    buffer.seek(0)
    return buffer


def process_latex_for_anki(text):
    """
    Converts LaTeX expressions to Anki's LaTeX format.
    """
    import re
    
    # Convert inline LaTeX ($...$) to Anki format
    text = re.sub(r'\$([^$]+)\$', r'\(\1\)', text)
    
    # Convert display LaTeX ($$...$$) to Anki format
    text = re.sub(r'\$\$([^$]+)\$\$', r'\[\1\]', text)
    
    return text


def export_flashcards_to_anki_apkg(flashcards):
    """
    Exports flashcards as a proper Anki .apkg file using genanki library.
    """
    try:
        import genanki
        from io import BytesIO
        import random
        from utils.gpt.shared import generate_ai_title
        
        # Generate AI title based on flashcard content
        content_sample = ""
        for i, (q, a) in enumerate(flashcards[:3]):  # Use first 3 cards for context
            content_sample += f"Q: {q}\nA: {a}\n"
        
        deck_title = generate_ai_title(content_sample, "flashcard deck")
        
        # Create a unique model ID for the card type
        model_id = random.randrange(1 << 30, 1 << 31)
        
        # Create a simple card model (front and back) with LaTeX support
        model = genanki.Model(
            model_id,
            'Simple Model',
            fields=[
                {'name': 'Question'},
                {'name': 'Answer'},
            ],
            templates=[
                {
                    'name': 'Card 1',
                    'qfmt': '{{Question}}',
                    'afmt': '{{FrontSide}}<hr id="answer">{{Answer}}',
                },
            ],
            css='''
            .card {
                font-family: arial;
                font-size: 20px;
                text-align: center;
                color: black;
                background-color: white;
            }
            '''
        )
        
        # Create a deck with AI-generated title
        deck_id = random.randrange(1 << 30, 1 << 31)
        deck = genanki.Deck(deck_id, deck_title)
        
        # Add cards to the deck with LaTeX processing
        for question, answer in flashcards:
            # Convert LaTeX to Anki format
            processed_question = process_latex_for_anki(question)
            processed_answer = process_latex_for_anki(answer)
            
            note = genanki.Note(
                model=model,
                fields=[processed_question, processed_answer]
            )
            deck.add_note(note)
        
        # Create the package
        package = genanki.Package(deck)
        
        # Generate the .apkg file in memory
        buffer = BytesIO()
        package.write_to_file(buffer)
        buffer.seek(0)
        
        return buffer
        
    except ImportError:
        # Fallback to TSV if genanki is not installed
        print("⚠️ genanki library not found. Installing...")
        import subprocess
        import sys
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "genanki"])
            # Retry after installation
            return export_flashcards_to_anki_apkg(flashcards)
        except:
            print("❌ Failed to install genanki. Falling back to TSV format.")
            return export_flashcards_to_anki_tsv(flashcards)
    except Exception as e:
        print(f"❌ Error creating Anki package: {e}. Falling back to TSV format.")
        return export_flashcards_to_anki_tsv(flashcards)
